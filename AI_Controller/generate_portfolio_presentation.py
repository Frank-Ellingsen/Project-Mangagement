import os
import datetime
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

def draw_gantt_chart(slide, x_offset, y_start, width, row_height, projects):
    # Timeframe Oct 1, 2025 to Dec 31, 2026
    start_date = datetime.date(2025, 10, 1)
    end_date = datetime.date(2026, 12, 31)
    total_days = (end_date - start_date).days
    
    # RAG Colors
    c_red = RGBColor(231, 76, 60)
    c_green = RGBColor(46, 204, 113)
    c_gray = RGBColor(203, 213, 225)
    
    # Draw timeline quarters
    quarters = [
        ("Q4 2025", datetime.date(2025, 10, 1), datetime.date(2025, 12, 31)),
        ("Q1 2026", datetime.date(2026, 1, 1), datetime.date(2026, 3, 31)),
        ("Q2 2026", datetime.date(2026, 4, 1), datetime.date(2026, 6, 30)),
        ("Q3 2026", datetime.date(2026, 7, 1), datetime.date(2026, 9, 30)),
        ("Q4 2026", datetime.date(2026, 10, 1), datetime.date(2026, 12, 31))
    ]
    
    # Draw timeline headers
    for name, q_start, q_end in quarters:
        q_days = (q_end - q_start).days + 1
        q_left = x_offset + width * ((q_start - start_date).days / total_days)
        q_width = width * (q_days / total_days)
        
        # Quarter Box
        q_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(q_left), Inches(y_start - 0.4), Inches(q_width), Inches(0.3))
        q_box.fill.solid()
        q_box.fill.fore_color.rgb = RGBColor(241, 245, 249)
        q_box.line.color.rgb = RGBColor(226, 232, 240)
        p = q_box.text_frame.paragraphs[0]
        p.text = name
        p.font.name = 'Inter'
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(100, 116, 139)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
    for i, proj in enumerate(projects):
        y_curr = y_start + i * row_height
        
        # Draw project label
        label_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_curr), Inches(4.0), Inches(0.5))
        p_lbl = label_box.text_frame.paragraphs[0]
        p_lbl.text = f"{proj['id']}: {proj['name']}"
        p_lbl.font.name = 'Outfit'
        p_lbl.font.size = Pt(14)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = RGBColor(26, 37, 47)
        
        # Draw baseline (Planned)
        p_start = datetime.datetime.strptime(proj["plannedStart"], "%Y-%m-%d").date()
        p_end = datetime.datetime.strptime(proj["plannedFinish"], "%Y-%m-%d").date()
        
        p_left = x_offset + width * ((p_start - start_date).days / total_days)
        p_width = width * (((p_end - p_start).days) / total_days)
        
        plan_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(p_left), Inches(y_curr + 0.05), Inches(p_width), Inches(0.1))
        plan_shape.fill.solid()
        plan_shape.fill.fore_color.rgb = c_gray
        plan_shape.line.fill.background()
        
        # Draw actual (Actual/Progress)
        a_start = datetime.datetime.strptime(proj["actualStart"], "%Y-%m-%d").date()
        a_end = datetime.datetime.strptime(proj["actualFinish"], "%Y-%m-%d").date()
        
        a_left = x_offset + width * ((a_start - start_date).days / total_days)
        a_width = width * (((a_end - a_start).days) / total_days)
        
        # Progress scale
        prog_width = a_width * (proj["progress"] / 100.0)
        
        actual_color = c_red if proj["status"] == "Red" else c_green
        bg_color = RGBColor(254, 226, 226) if proj["status"] == "Red" else RGBColor(220, 252, 231)
        
        # Actual background bar (muted status color)
        actual_shape_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(a_left), Inches(y_curr + 0.2), Inches(a_width), Inches(0.18))
        actual_shape_bg.fill.solid()
        actual_shape_bg.fill.fore_color.rgb = bg_color
        actual_shape_bg.line.fill.background()
        
        if prog_width > 0:
            actual_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(a_left), Inches(y_curr + 0.2), Inches(prog_width), Inches(0.18))
            actual_shape.fill.solid()
            actual_shape.fill.fore_color.rgb = actual_color
            actual_shape.line.fill.background()
            
        # Draw progress & RAG status text
        status_box = slide.shapes.add_textbox(Inches(x_offset + width + 0.1), Inches(y_curr + 0.1), Inches(2.0), Inches(0.4))
        p_stat = status_box.text_frame.paragraphs[0]
        status_text = "🔴 RED" if proj["status"] == "Red" else "🟢 GRN"
        p_stat.text = f"{proj['progress']:.1f}% | {status_text}"
        p_stat.font.name = 'Inter'
        p_stat.font.size = Pt(12)
        p_stat.font.bold = True
        p_stat.font.color.rgb = actual_color

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 10 slides data conforming to the 10-20-30 rule (10 slides, min 30pt font size)
    slides_data = [
        {
            "title": "Project Portfolio Executive Review",
            "subtitle": "Q2 2026 Status Report | Frank Ellingsen",
            "bullets": [
                "Comprehensive status of 6 active & completed projects",
                "EVM metrics overview & financial audits",
                "Strategic recommendations & corrective actions"
            ]
        },
        # Slide 2 will be the custom Gantt of Gantts slide
        {
            "title": "1. Portfolio Overview & Health",
            "bullets": [
                "Total Baseline (BAC): 7,100,000 USD",
                "Total Spent (AC): 5,220,810 USD",
                "Earned Value (EV): 4,752,500 USD",
                "Portfolio CPI: 0.91 (Muted cost overrun)"
            ]
        },
        {
            "title": "2. PRJ-001 Vessel Construction",
            "status": "Red",
            "bullets": [
                "Status: Completed (99.5% progress) | 🔴 RED",
                "Baseline: 1.5M USD | Actual Cost: 1.85M USD",
                "CPI: 0.81 | Total Overrun: -356K USD",
                "Action: Renegotiate contractor rates & freeze VOs"
            ]
        },
        {
            "title": "3. PRJ-002 Patrol Vessel Design",
            "status": "Green",
            "bullets": [
                "Status: Planned (0% progress) | 🟢 GREEN",
                "Baseline Budget: 800,000 USD",
                "Planned Start: Aug 2026 | Finish: Dec 2026",
                "Focus: Carbon mold design and CFD analysis"
            ]
        },
        {
            "title": "4. PRJ-003 Subsea Frame",
            "status": "Green",
            "bullets": [
                "Status: Active (30% progress) | 🟢 GREEN",
                "Baseline: 1.2M USD | Spent: 382.5K USD",
                "CPI: 0.94 (Tracking close to budget)",
                "Focus: Engineering complete, steel fabrication starting"
            ]
        },
        {
            "title": "5. PRJ-004 Workboat Hull",
            "status": "Green",
            "bullets": [
                "Status: Active (70% progress) | 🟢 GREEN",
                "Baseline: 2.0M USD | Spent: 1.48M USD",
                "CPI: 0.95 (Slight budget pressure)",
                "Focus: Drawings approved, welding & assembly ongoing"
            ]
        },
        {
            "title": "6. PRJ-005 Logistics Pontoon",
            "status": "Green",
            "bullets": [
                "Status: Active (90% progress) | 🟢 GREEN",
                "Baseline: 1.0M USD | Spent: 927K USD",
                "CPI: 0.97 (In good financial standing)",
                "Project Manager: Frank Ellingsen"
            ]
        },
        {
            "title": "7. PRJ-006 Cargo Hatch",
            "status": "Green",
            "bullets": [
                "Status: Completed (100% progress) | 🟢 GREEN",
                "Baseline: 600,000 USD | Spent: 586,500 USD",
                "CPI: 1.02 (Under budget success)",
                "Deliverables: Molding & testing fully verified"
            ]
        },
        {
            "title": "8. Key Portfolio Risks & RAID Log",
            "bullets": [
                "Risk: Carbon sheet delivery delay from supplier",
                "Issue: Welder overtime limits exceeded in fabrication",
                "Dependency: Sea trials depend on DNV class approval",
                "Mitigation: Allocate junior welders to assist"
            ]
        },
        {
            "title": "9. Recommendations & Next Steps",
            "bullets": [
                "Renegotiate outfitting contractor rates (WBS 1 & 3)",
                "Shift composite fabricators to outfitting to optimize",
                "Enforce double-shift schedules to recover sea trials",
                "Integrate DuckDB tracking into future estimations"
            ]
        }
    ]
    
    # Portfolio projects data for the Gantt Chart
    portfolio_projects = [
        {"id": "PRJ-001", "name": "Composite Vessel", "plannedStart": "2026-01-01", "plannedFinish": "2026-06-30", "actualStart": "2026-01-01", "actualFinish": "2026-06-30", "progress": 99.5, "status": "Red"},
        {"id": "PRJ-002", "name": "Patrol Vessel", "plannedStart": "2026-08-01", "plannedFinish": "2026-12-31", "actualStart": "2026-08-01", "actualFinish": "2026-12-31", "progress": 0.0, "status": "Green"},
        {"id": "PRJ-003", "name": "Subsea Frame", "plannedStart": "2026-05-01", "plannedFinish": "2026-10-31", "actualStart": "2026-05-01", "actualFinish": "2026-10-31", "progress": 30.0, "status": "Green"},
        {"id": "PRJ-004", "name": "Workboat Hull", "plannedStart": "2026-03-01", "plannedFinish": "2026-08-31", "actualStart": "2026-03-01", "actualFinish": "2026-08-31", "progress": 70.0, "status": "Green"},
        {"id": "PRJ-005", "name": "Logistics Pontoon", "plannedStart": "2026-02-01", "plannedFinish": "2026-07-31", "actualStart": "2026-02-01", "actualFinish": "2026-07-31", "progress": 90.0, "status": "Green"},
        {"id": "PRJ-006", "name": "Cargo Hatch", "plannedStart": "2025-10-01", "plannedFinish": "2026-03-31", "actualStart": "2025-10-01", "actualFinish": "2026-03-31", "progress": 100.0, "status": "Green"}
    ]
    
    # Muted clean Tufte colors
    navy = RGBColor(26, 37, 47)
    slate = RGBColor(44, 62, 80)
    dark_gray = RGBColor(85, 85, 85)
    
    # RAG Highlight colors
    red_color = RGBColor(231, 76, 60)
    green_color = RGBColor(46, 204, 113)
    
    for idx, slide_info in enumerate(slides_data):
        # Insert Title Slide
        if idx == 0:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Title Box
            title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(1.5))
            tf_title = title_box.text_frame
            tf_title.word_wrap = True
            p_title = tf_title.paragraphs[0]
            p_title.text = slide_info["title"]
            p_title.font.name = 'Outfit'
            p_title.font.size = Pt(44)
            p_title.font.bold = True
            p_title.font.color.rgb = navy
            
            # Subtitle & Bullets Content
            content_box = slide.shapes.add_textbox(Inches(0.75), Inches(3.2), Inches(11.83), Inches(3.5))
            tf_content = content_box.text_frame
            tf_content.word_wrap = True
            
            p_sub = tf_content.paragraphs[0]
            p_sub.text = slide_info["subtitle"]
            p_sub.font.name = 'Inter'
            p_sub.font.size = Pt(30)
            p_sub.font.bold = True
            p_sub.font.color.rgb = slate
            p_sub.space_after = Pt(20)
            
            for b in slide_info["bullets"]:
                p = tf_content.add_paragraph()
                p.text = "• " + b
                p.font.name = 'Inter'
                p.font.size = Pt(30)
                p.font.color.rgb = dark_gray
                p.space_after = Pt(14)
                
            # Add custom Gantt of Gantts slide right after the Title slide
            gantt_slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            g_title_box = gantt_slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(1.0))
            tf_g_title = g_title_box.text_frame
            p_g_title = tf_g_title.paragraphs[0]
            p_g_title.text = "Portfolio Gantt of Gantts"
            p_g_title.font.name = 'Outfit'
            p_g_title.font.size = Pt(40)
            p_g_title.font.bold = True
            p_g_title.font.color.rgb = navy
            
            # Draw Gantt
            draw_gantt_chart(gantt_slide, x_offset=4.5, y_start=1.8, width=6.2, row_height=0.8, projects=portfolio_projects)
            continue
            
        # Standard content slides
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title Box
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(1.5))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = slide_info["title"]
        p_title.font.name = 'Outfit'
        p_title.font.size = Pt(44)
        p_title.font.bold = True
        p_title.font.color.rgb = navy
        
        # If RAG status is defined, add color highlighting to the title or a badge
        if "status" in slide_info:
            badge_color = red_color if slide_info["status"] == "Red" else green_color
            status_text = "🔴 CRITICAL" if slide_info["status"] == "Red" else "🟢 ON TRACK"
            
            badge_box = slide.shapes.add_textbox(Inches(9.5), Inches(0.65), Inches(3.0), Inches(0.8))
            p_badge = badge_box.text_frame.paragraphs[0]
            p_badge.text = status_text
            p_badge.font.name = 'Inter'
            p_badge.font.size = Pt(20)
            p_badge.font.bold = True
            p_badge.font.color.rgb = badge_color
            
        # Content Box
        content_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(11.83), Inches(4.5))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        
        first = True
        for b in slide_info["bullets"]:
            if first:
                p = tf_content.paragraphs[0]
                first = False
            else:
                p = tf_content.add_paragraph()
            p.text = "• " + b
            p.font.name = 'Inter'
            p.font.size = Pt(30)
            
            # Apply RAG color to specific bullets containing status terms
            if "🔴 RED" in b or "Red" in b:
                p.font.color.rgb = red_color
            elif "🟢 GREEN" in b or "Green" in b or "Success" in b:
                p.font.color.rgb = green_color
            else:
                p.font.color.rgb = dark_gray
                
            p.space_after = Pt(20)
                
    # Save presentation
    os.makedirs("Data", exist_ok=True)
    prs.save("Data/Project_Portfolio_Review.pptx")
    print("PowerPoint presentation with Gantt of Gantts generated successfully!")

if __name__ == "__main__":
    create_presentation()
